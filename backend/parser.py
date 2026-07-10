import io
import pdfplumber
from docx import Document
from pptx import Presentation

def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extracts text from PDF bytes using pdfplumber."""
    text_content = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_content.append(page_text)
    return "\n".join(text_content)

def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extracts text from DOCX bytes using python-docx."""
    doc = Document(io.BytesIO(file_bytes))
    text_content = []
    for paragraph in doc.paragraphs:
        if paragraph.text:
            text_content.append(paragraph.text)
    
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text for cell in row.cells if cell.text]
            if row_text:
                text_content.append(" | ".join(row_text))
                
    return "\n".join(text_content)

def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extracts text from PPTX bytes using python-pptx."""
    prs = Presentation(io.BytesIO(file_bytes))
    text_content = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        slide_text.append(paragraph.text)
        if slide_text:
            text_content.append("\n".join(slide_text))
    return "\n\n".join(text_content)

def extract_text(file_bytes: bytes, file_name: str) -> str:
    """Detects file type by extension and extracts text."""
    lower_name = file_name.lower()
    if lower_name.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    elif lower_name.endswith(".docx"):
        return extract_text_from_docx(file_bytes)
    elif lower_name.endswith(".pptx") or lower_name.endswith(".ppt"):
        return extract_text_from_pptx(file_bytes)
    else:
        raise ValueError("Format file tidak didukung. Gunakan PDF, DOCX, atau PPTX.")
